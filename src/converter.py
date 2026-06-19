"""V2.6.3.10 — 비-마크다운 → 마크다운 변환기.
V2.6.3.11 — 스캔 PDF OCR fallback 추가.

목적: PDF/PPTX/DOCX 등 비-텍스트 포맷도 archive 우선 파이프라인에 박을 수 있게.
      original.<ext>는 그대로 보존, content.md는 변환된 마크다운.

지원 포맷:
  .md, .txt     → passthrough (frontmatter 제거)
  .pdf          → pdfplumber 추출 (페이지 단위), 스캔본 의심 시 Tesseract OCR fallback
  .pptx, .ppt   → python-pptx (슬라이드 단위)
  .docx, .doc   → mammoth (마크다운 직접 변환)

OCR 정책 (V2.6.3.11):
  - PDF 추출 결과의 페이지당 평균 글자수 < IRIS_OCR_MIN_CHARS (기본 20)
    또는 전체 결과가 비어 있으면 → OCR fallback
  - 언어: IRIS_OCR_LANG 환경변수, 기본 "kor+chi_sim+eng"
  - Tesseract 미설치 시: 텍스트 결과만 반환 (OCR 안 함)

extraction_method (호출자에 hint):
  - "text"  : pdfplumber 결과만 사용
  - "ocr"   : 텍스트 빈 페이지를 OCR로 채움
  - "hybrid": 일부 텍스트 + 일부 OCR (현 시점 미사용, 미래 확장)
"""
from __future__ import annotations

import os
from pathlib import Path


SUPPORTED_SUFFIXES = {".md", ".txt", ".pdf", ".pptx", ".ppt", ".docx", ".doc"}
LOSSLESS_SUFFIXES = {".md", ".txt"}   # 마크다운/텍스트 = original==content


class ConversionError(Exception):
    """변환 실패. 메시지에 원인 박힘."""


class ConversionResult:
    """변환 결과 + 메타. 호출자가 manifest에 박을 수 있게."""
    __slots__ = ("body", "extraction_method", "pages", "ocr_pages")

    def __init__(self, body: str, extraction_method: str = "text",
                 pages: int = 0, ocr_pages: int = 0):
        self.body = body
        self.extraction_method = extraction_method  # "text" | "ocr" | "hybrid"
        self.pages = pages
        self.ocr_pages = ocr_pages


def convert_to_markdown(src: Path) -> str:
    """src 파일을 마크다운 문자열로 변환 (간단 API — 본문만 반환).

    raise ConversionError: 라이브러리 부재, 파싱 실패, 파일 없음.
    빈 결과 ("") 정상 — 호출자가 skip-empty.
    """
    return convert_with_meta(src).body


def convert_with_meta(src: Path) -> ConversionResult:
    """변환 + 메타 반환. PDF의 OCR 사용 여부 등 호출자가 manifest에 박을 수 있게."""
    if not src.exists():
        raise ConversionError(f"파일 없음: {src}")

    suf = src.suffix.lower()
    if suf in {".md", ".txt"}:
        return ConversionResult(_convert_text(src), extraction_method="text")
    elif suf == ".pdf":
        return _convert_pdf_with_meta(src)
    elif suf in {".pptx", ".ppt"}:
        return ConversionResult(_convert_pptx(src), extraction_method="text")
    elif suf in {".docx", ".doc"}:
        return ConversionResult(_convert_docx(src), extraction_method="text")
    else:
        raise ConversionError(f"지원 안 함: {suf}")


# ─── 포맷별 변환기 ────────────────────────────────────────────────────
def _convert_text(src: Path) -> str:
    """마크다운/텍스트 — frontmatter만 제거."""
    text = src.read_text(encoding="utf-8")
    if text.startswith("---\n"):
        end = text.find("\n---\n", 4)
        if end > 0:
            return text[end + 5:]
        end = text.find("\n---", 4)
        if end > 0:
            return text[end + 4:].lstrip("\n")
    return text


def _convert_pdf_with_meta(src: Path) -> ConversionResult:
    """PDF — pdfplumber로 페이지 추출. 스캔본 의심 시 OCR fallback (V2.6.3.11)."""
    try:
        import pdfplumber
    except ImportError as e:
        raise ConversionError(f"pdfplumber 미설치: {e}")

    OCR_MIN_CHARS = int(os.environ.get("IRIS_OCR_MIN_CHARS", "20"))

    out: list[str] = []
    text_total_chars = 0
    page_count = 0
    page_texts: list[tuple[int, str, list]] = []  # (page_no, text, tables)

    try:
        with pdfplumber.open(src) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                text = (page.extract_text() or "").strip()
                tables = page.extract_tables() or []
                page_texts.append((i, text, tables))
                text_total_chars += len(text)
    except Exception as e:
        raise ConversionError(f"PDF 파싱 실패: {type(e).__name__}: {e}")

    # OCR fallback 판단 — 페이지당 평균 < OCR_MIN_CHARS
    avg_chars = text_total_chars / page_count if page_count else 0
    needs_ocr = page_count > 0 and avg_chars < OCR_MIN_CHARS

    extraction_method = "text"
    ocr_pages = 0

    if needs_ocr:
        ocr_results = _ocr_pdf(src, page_count)
        if ocr_results is not None:
            # OCR 결과로 빈 페이지 보충
            extraction_method = "ocr"
            for i, ocr_text in ocr_results:
                # 같은 페이지 인덱스의 텍스트가 비어 있으면 OCR로 교체
                for j, (pno, text, tables) in enumerate(page_texts):
                    if pno == i and len(text) < OCR_MIN_CHARS:
                        page_texts[j] = (pno, ocr_text.strip(), tables)
                        ocr_pages += 1
                        break

    # 본문 조립
    for pno, text, tables in page_texts:
        page_md = []
        if text:
            page_md.append(text)
        for t in tables:
            md_table = _list_table_to_markdown(t)
            if md_table:
                page_md.append(md_table)
        if page_md:
            out.append(f"## 페이지 {pno}\n\n" + "\n\n".join(page_md))

    return ConversionResult(
        body="\n\n".join(out),
        extraction_method=extraction_method,
        pages=page_count,
        ocr_pages=ocr_pages,
    )


def _ocr_pdf(src: Path, page_count: int) -> list[tuple[int, str]] | None:
    """PDF → 이미지 → Tesseract OCR. 페이지별 (페이지번호, 텍스트) 리스트.

    None 반환 = Tesseract/pdf2image 미설치 등으로 OCR 자체 불가.
    """
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except ImportError:
        return None

    lang = os.environ.get("IRIS_OCR_LANG", "kor+chi_sim+eng")
    dpi = int(os.environ.get("IRIS_OCR_DPI", "200"))

    try:
        images = convert_from_path(str(src), dpi=dpi)
    except Exception:
        # poppler 미설치 또는 PDF 손상
        return None

    results: list[tuple[int, str]] = []
    for i, img in enumerate(images, 1):
        try:
            text = pytesseract.image_to_string(img, lang=lang)
            results.append((i, text))
        except Exception:
            results.append((i, ""))
    return results


# 하위 호환 — 옛 시그니처 유지
def _convert_pdf(src: Path) -> str:
    return _convert_pdf_with_meta(src).body


def _convert_pptx(src: Path) -> str:
    """PPTX — python-pptx로 슬라이드 단위 추출."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ConversionError(f"python-pptx 미설치: {e}")

    out: list[str] = []
    try:
        prs = Presentation(str(src))
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = [f"## 슬라이드 {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join(run.text for run in para.runs).strip()
                        if txt:
                            slide_lines.append(txt)
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in tbl.rows
                    ]
                    md_table = _list_table_to_markdown(rows)
                    if md_table:
                        slide_lines.append(md_table)
            # 노트
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_lines.append(f"\n> **노트**: {notes}")
            except Exception:
                pass
            if len(slide_lines) > 1:
                out.append("\n\n".join(slide_lines))
    except Exception as e:
        raise ConversionError(f"PPTX 파싱 실패: {type(e).__name__}: {e}")
    return "\n\n".join(out)


def _convert_docx(src: Path) -> str:
    """DOCX — mammoth로 마크다운 직접 변환."""
    try:
        import mammoth
    except ImportError as e:
        raise ConversionError(f"mammoth 미설치: {e}")

    try:
        with src.open("rb") as f:
            result = mammoth.convert_to_markdown(f)
        return result.value or ""
    except Exception as e:
        raise ConversionError(f"DOCX 파싱 실패: {type(e).__name__}: {e}")


# ─── 유틸 ────────────────────────────────────────────────────────────
def _list_table_to_markdown(rows: list[list]) -> str:
    """2D 리스트 → 마크다운 표. None/빈 셀은 공백으로."""
    if not rows:
        return ""
    # 행 정리
    cleaned = [
        [(str(c) if c is not None else "").replace("\n", " ").strip() for c in row]
        for row in rows
    ]
    cleaned = [r for r in cleaned if any(r)]  # 빈 줄 제거
    if not cleaned:
        return ""
    ncols = max(len(r) for r in cleaned)
    # 패딩
    for r in cleaned:
        while len(r) < ncols:
            r.append("")
    header = cleaned[0]
    out = ["| " + " | ".join(header) + " |"]
    out.append("| " + " | ".join("---" for _ in header) + " |")
    for r in cleaned[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


__all__ = [
    "SUPPORTED_SUFFIXES", "LOSSLESS_SUFFIXES",
    "ConversionError", "ConversionResult",
    "convert_to_markdown", "convert_with_meta",
]
