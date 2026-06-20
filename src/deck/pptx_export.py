"""V2.7.6 — Deck → PPTX (PNG 임베드)."""
from __future__ import annotations

import tempfile
from pathlib import Path

from src.deck.schema import Deck


def render_deck_to_pptx(deck: Deck, out_path: Path | None = None) -> Path:
    from pptx import Presentation
    from pptx.util import Inches
    from playwright.sync_api import sync_playwright

    from src.deck.renderer import render_slide_html

    if out_path is None:
        out_path = Path(tempfile.mkdtemp(prefix="iris_pptx_")) / "deck.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)

    tmp_dir = Path(tempfile.mkdtemp(prefix="iris_deck_png_"))
    png_paths: list[Path] = []
    total = len(deck.slides)

    with sync_playwright() as p:
        browser = p.chromium.launch()
        ctx = browser.new_context(viewport={"width": 1920, "height": 1080})
        page = ctx.new_page()
        for i, slide in enumerate(deck.slides, 1):
            html = render_slide_html(
                slide.pattern, slide.data, deck,
                pageno=i, total_pages=total,
            )
            html_path = tmp_dir / f"slide_{i:03d}.html"
            html_path.write_text(html, encoding="utf-8")
            page.goto(f"file://{html_path.resolve()}")
            png_path = tmp_dir / f"slide_{i:03d}.png"
            page.screenshot(path=str(png_path), full_page=True)
            png_paths.append(png_path)
        browser.close()

    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    blank = prs.slide_layouts[6]
    for png in png_paths:
        sl = prs.slides.add_slide(blank)
        sl.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(out_path))
    return out_path


__all__ = ["render_deck_to_pptx"]
