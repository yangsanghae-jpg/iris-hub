"""V2.6.3.10 — 비-마크다운 → 마크다운 변환기.
V2.6.3.11 — 스캔 PDF OCR fallback 추가.
V2.9 — PPT 탭 소스 업로드: .xlsx 지원 추가, enable_ocr 플래그.

목적: PDF/PPTX/DOCX/XLSX 등 비-텍스트 포맷도 archive 우선 파이프라인에 박을 수 있게.
      original.<ext>는 그대로 보존, content.md는 변환된 마크다운.

지원 포맷:
  .md, .txt     → passthrough (frontmatter 제거)
  .pdf          → pdfplumber 추출 (페이지 단위), 스캔본 의심 시 Tesseract OCR fallback
  .pptx, .ppt   → python-pptx (슬라이드 단위)
  .docx, .doc   → mammoth (마크다운 직접 변환)
  .xlsx         → openpyxl (시트 단위 표)

레거시 바이너리 포맷(.ppt/.doc/.xls)은 converter 자체는 계속 지원하지만,
PPT 탭 업로드 API(v31/server.py)는 OOXML만 받는다 — LEGACY_BINARY_SUFFIXES 참조.

OCR 정책 (V2.6.3.11):
  - PDF 추출 결과의 페이지당 평균 글자수 < IRIS_OCR_MIN_CHARS (기본 20)
    또는 전체 결과가 비어 있으면 → OCR fallback
  - 언어: IRIS_OCR_LANG 환경변수, 기본 "kor+chi_sim+eng"
  - Tesseract 미설치 시: 텍스트 결과만 반환 (OCR 안 함)

extraction_method (호출자에 hint):
  - "text"  : pdfplumber 결과만 사용
  - "ocr"   : 텍스트 빈 페이지를 OCR로 채움
  - "hybrid": 일부 텍스트 + 일부 OCR (현 시점 미사용, 미래 확장)
"""
from __future__ import annotations

import os
from pathlib import Path


SUPPORTED_SUFFIXES = {".md", ".txt", ".pdf", ".pptx", ".ppt", ".docx", ".doc", ".xlsx"}
LOSSLESS_SUFFIXES = {".md", ".txt"}   # 마크다운/텍스트 = original==content

# PPT 탭 업로드 등에서 "레거시 바이너리 → OOXML로 저장 요청" 안내에 사용.
# converter 자체는 .ppt/.doc를 계속 변환하지만, 신규 업로드 경로는 거부한다.
LEGACY_BINARY_SUFFIXES = {".ppt", ".doc", ".xls"}


class ConversionError(Exception):
    """변환 실패. 메시지에 원인 박힘."""


class ConversionResult:
    """변환 결과 + 메타. 호출자가 manifest에 박을 수 있게."""
    __slots__ = ("body", "extraction_method", "pages", "ocr_pages")

    def __init__(self, body: str, extraction_method: str = "text",
                 pages: int = 0, ocr_pages: int = 0):
        self.body = body
        self.extraction_method = extraction_method  # "text" | "ocr" | "hybrid"
        self.pages = pages
        self.ocr_pages = ocr_pages


def convert_to_markdown(src: Path) -> str:
    """src 파일을 마크다운 문자열로 변환 (간단 API — 본문만 반환).

    raise ConversionError: 라이브러리 부재, 파싱 실패, 파일 없음.
    빈 결과 ("") 정상 — 호출자가 skip-empty.
    """
    return convert_with_meta(src).body


def convert_with_meta(src: Path, enable_ocr: bool = True) -> ConversionResult:
    """변환 + 메타 반환. PDF의 OCR 사용 여부 등 호출자가 manifest에 박을 수 있게.

    enable_ocr=False면 PDF 변환 시 OCR fallback을 건너뛰고 텍스트 추출만 사용한다
    (빠른 미리보기 등). 기본값 True — 기존 호출자(folder_load 등) 동작 보존.
    """
    if not src.exists():
        raise ConversionError(f"파일 없음: {src}")

    suf = src.suffix.lower()
    if suf in {".md", ".txt"}:
        return ConversionResult(_convert_text(src), extraction_method="text")
    elif suf == ".pdf":
        return _convert_pdf_with_meta(src, enable_ocr=enable_ocr)
    elif suf in {".pptx", ".ppt"}:
        return ConversionResult(_convert_pptx(src), extraction_method="text")
    elif suf in {".docx", ".doc"}:
        return ConversionResult(_convert_docx(src), extraction_method="text")
    elif suf == ".xlsx":
        return ConversionResult(_convert_xlsx(src), extraction_method="text")
    else:
        raise ConversionError(f"지원 안 함: {suf}")


# ─── 포맷별 변환기 ────────────────────────────────────────────────────
def _convert_text(src: Path) -> str:
    """마크다운/텍스트 — frontmatter만 제거."""
    text = src.read_text(encoding="utf-8")
    if text.startswith("---\n"):
        end = text.find("\n---\n", 4)
        if end > 0:
            return text[end + 5:]
        end = text.find("\n---", 4)
        if end > 0:
            return text[end + 4:].lstrip("\n")
    return text


def _convert_pdf_with_meta(src: Path, enable_ocr: bool = True) -> ConversionResult:
    """PDF — 페이지별 분기 (V2.7.4 hybrid).

    페이지마다:
      - pdfplumber 텍스트 추출 → IRIS_OCR_MIN_CHARS 미만이면 OCR fallback
      - 일부 페이지만 OCR하는 경우 extraction_method = "hybrid"

    이전 V2.6.3.11은 *전체 평균*으로 판단 → 부분 스캔 PDF에서 텍스트 페이지 손실
    또는 불필요한 전체 OCR 발생. V2.7.4는 페이지별 결정으로 정밀화.

    enable_ocr=False면 OCR fallback을 건너뛰고 텍스트 추출 결과만 사용한다.
    """
    try:
        import pdfplumber
    except ImportError as e:
        raise ConversionError(f"pdfplumber 미설치: {e}")

    OCR_MIN_CHARS = int(os.environ.get("IRIS_OCR_MIN_CHARS", "20"))
    OCR_MAX_PAGES = int(os.environ.get("IRIS_OCR_MAX_PAGES", "200"))  # 안전상한

    # 1) 모든 페이지 텍스트 + 표 한 번에 추출
    page_data: list[tuple[int, str, list]] = []  # (page_no, text, tables)
    try:
        with pdfplumber.open(src) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                text = (page.extract_text() or "").strip()
                tables = page.extract_tables() or []
                page_data.append((i, text, tables))
    except Exception as e:
        raise ConversionError(f"PDF 파싱 실패: {type(e).__name__}: {e}")

    if page_count == 0:
        return ConversionResult(body="", extraction_method="text", pages=0)

    # 2) OCR 필요한 페이지 번호 식별 (enable_ocr=False면 전부 스킵)
    ocr_needed_pages = [
        pno for pno, text, _tables in page_data
        if len(text) < OCR_MIN_CHARS
    ] if enable_ocr else []

    # 3) OCR 일괄 수행 (필요한 페이지만)
    ocr_pages = 0
    ocr_results: dict[int, str] = {}
    if ocr_needed_pages:
        if len(ocr_needed_pages) > OCR_MAX_PAGES:
            # 안전 상한 — 앞 OCR_MAX_PAGES만
            ocr_needed_pages = ocr_needed_pages[:OCR_MAX_PAGES]
        ocr_results = _ocr_pdf_pages(src, ocr_needed_pages)
        ocr_pages = sum(1 for pno in ocr_needed_pages if ocr_results.get(pno))

    # 4) 페이지 데이터에 OCR 결과 머지
    merged: list[tuple[int, str, list]] = []
    text_pages_used = 0
    for pno, text, tables in page_data:
        if pno in ocr_results and ocr_results[pno]:
            merged.append((pno, ocr_results[pno].strip(), tables))
        else:
            merged.append((pno, text, tables))
            if len(text) >= OCR_MIN_CHARS:
                text_pages_used += 1

    # 5) extraction_method 결정
    if ocr_pages == 0:
        extraction_method = "text"
    elif text_pages_used == 0:
        extraction_method = "ocr"
    else:
        extraction_method = "hybrid"

    # 6) 본문 조립 + 페이지 표식 (어떤 방식으로 추출됐는지 미세 표시)
    out: list[str] = []
    for pno, text, tables in merged:
        page_md = []
        if text:
            page_md.append(text)
        for t in tables:
            md_table = _list_table_to_markdown(t)
            if md_table:
                page_md.append(md_table)
        if page_md:
            # hybrid일 때 OCR 페이지에는 작은 표식 (디버그용)
            marker = " · OCR" if pno in ocr_results and ocr_results.get(pno) else ""
            out.append(f"## 페이지 {pno}{marker}\n\n" + "\n\n".join(page_md))

    return ConversionResult(
        body="\n\n".join(out),
        extraction_method=extraction_method,
        pages=page_count,
        ocr_pages=ocr_pages,
    )


def _ocr_pdf_pages(src: Path, pages: list[int]) -> dict[int, str]:
    """지정한 페이지 번호들만 OCR. V2.7.4 — pdf2image first_page/last_page 활용.

    반환: {page_no: text}. 라이브러리 미설치 시 빈 dict.
    """
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except ImportError:
        return {}

    if not pages:
        return {}

    lang = os.environ.get("IRIS_OCR_LANG", "kor+chi_sim+eng")
    dpi = int(os.environ.get("IRIS_OCR_DPI", "200"))

    results: dict[int, str] = {}
    # pdf2image는 연속 범위가 효율적. 연속된 페이지 그룹화.
    pages_sorted = sorted(pages)
    ranges: list[tuple[int, int]] = []
    start = pages_sorted[0]
    prev = start
    for p in pages_sorted[1:]:
        if p == prev + 1:
            prev = p
        else:
            ranges.append((start, prev))
            start = prev = p
    ranges.append((start, prev))

    for first, last in ranges:
        try:
            images = convert_from_path(
                str(src), dpi=dpi,
                first_page=first, last_page=last,
            )
        except Exception:
            continue
        for offset, img in enumerate(images):
            pno = first + offset
            try:
                text = pytesseract.image_to_string(img, lang=lang)
                results[pno] = text
            except Exception:
                results[pno] = ""
    return results


# 하위 호환 — 옛 시그니처 유지 (V2.6.3.11 이전 호출자)
def _ocr_pdf(src: Path, page_count: int) -> list[tuple[int, str]] | None:
    """전체 페이지 OCR. V2.7.4부터 권장 안 함 (_ocr_pdf_pages 사용)."""
    all_pages = list(range(1, page_count + 1))
    results = _ocr_pdf_pages(src, all_pages)
    if not results:
        return None
    return [(p, results.get(p, "")) for p in all_pages]


# 하위 호환 — 옛 시그니처 유지
def _convert_pdf(src: Path) -> str:
    return _convert_pdf_with_meta(src).body


def _convert_pptx(src: Path) -> str:
    """PPTX — python-pptx로 슬라이드 단위 추출."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ConversionError(f"python-pptx 미설치: {e}")

    out: list[str] = []
    try:
        prs = Presentation(str(src))
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = [f"## 슬라이드 {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join(run.text for run in para.runs).strip()
                        if txt:
                            slide_lines.append(txt)
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in tbl.rows
                    ]
                    md_table = _list_table_to_markdown(rows)
                    if md_table:
                        slide_lines.append(md_table)
            # 노트
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_lines.append(f"\n> **노트**: {notes}")
            except Exception:
                pass
            if len(slide_lines) > 1:
                out.append("\n\n".join(slide_lines))
    except Exception as e:
        raise ConversionError(f"PPTX 파싱 실패: {type(e).__name__}: {e}")
    return "\n\n".join(out)


def _convert_docx(src: Path) -> str:
    """DOCX — mammoth로 마크다운 직접 변환."""
    try:
        import mammoth
    except ImportError as e:
        raise ConversionError(f"mammoth 미설치: {e}")

    try:
        with src.open("rb") as f:
            result = mammoth.convert_to_markdown(f)
        return result.value or ""
    except Exception as e:
        raise ConversionError(f"DOCX 파싱 실패: {type(e).__name__}: {e}")


def _convert_xlsx(src: Path) -> str:
    """XLSX — openpyxl로 시트 단위 표 추출 (수식은 계산된 값으로)."""
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ConversionError(f"openpyxl 미설치: {e}")

    out: list[str] = []
    try:
        wb = load_workbook(filename=str(src), read_only=True, data_only=True)
        try:
            for name in wb.sheetnames:
                ws = wb[name]
                rows = [list(row) for row in ws.iter_rows(values_only=True)]
                md_table = _list_table_to_markdown(rows)
                if md_table:
                    out.append(f"## 시트: {name}\n\n{md_table}")
        finally:
            wb.close()
    except ConversionError:
        raise
    except Exception as e:
        raise ConversionError(f"XLSX 파싱 실패: {type(e).__name__}: {e}")
    return "\n\n".join(out)


# ─── 유틸 ────────────────────────────────────────────────────────────
def _list_table_to_markdown(rows: list[list]) -> str:
    """2D 리스트 → 마크다운 표. None/빈 셀은 공백으로."""
    if not rows:
        return ""
    # 행 정리
    cleaned = [
        [(str(c) if c is not None else "").replace("\n", " ").strip() for c in row]
        for row in rows
    ]
    cleaned = [r for r in cleaned if any(r)]  # 빈 줄 제거
    if not cleaned:
        return ""
    ncols = max(len(r) for r in cleaned)
    # 패딩
    for r in cleaned:
        while len(r) < ncols:
            r.append("")
    header = cleaned[0]
    out = ["| " + " | ".join(header) + " |"]
    out.append("| " + " | ".join("---" for _ in header) + " |")
    for r in cleaned[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


__all__ = [
    "SUPPORTED_SUFFIXES", "LOSSLESS_SUFFIXES", "LEGACY_BINARY_SUFFIXES",
    "ConversionError", "ConversionResult",
    "convert_to_markdown", "convert_with_meta",
]
