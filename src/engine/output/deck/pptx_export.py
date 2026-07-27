"""V2.8.0 — Deck → PPTX (네이티브 도형/텍스트박스, PowerPoint에서 직접 편집 가능)."""
from __future__ import annotations

import math
import tempfile
from contextvars import ContextVar
from pathlib import Path
from typing import TYPE_CHECKING

from src.engine.output.deck.schema import Deck, detect_deck_lang

if TYPE_CHECKING:
    from src.engine.output.deck.theme import ResolvedDeckStyle

SLIDE_W_IN = 20.0
SLIDE_H_IN = 11.25
MARGIN = 0.625
HEADER_H = 1.1
CONTENT_TOP = 1.5
FOOTER_TOP = 10.75
CONTENT_BOTTOM = 10.55
CONTENT_W = SLIDE_W_IN - 2 * MARGIN

_FONT_BY_LANG = {"ko": "Malgun Gothic", "zh": "Microsoft YaHei", "en": "Segoe UI"}
_COVER_LABELS = {
    "ko": ("보고 대상", "보고 일자 / 버전", "Confidential — 사내 한정 사용"),
    "zh": ("报告对象", "报告日期 / 版本", "Confidential — 内部使用"),
    "en": ("Target Audience", "Date / Version", "Confidential — Internal Use Only"),
}
# 렌더 시작 시 _detect_lang()로 한 번 정해서 여기 담아둠 — 전체 슬라이드가 같은
# 폰트/라벨 언어를 쓰게. 모듈 전역이라 render_deck_to_pptx 호출마다 재설정됨.
FONT = _FONT_BY_LANG["ko"]
_cover_labels = _COVER_LABELS["ko"]

# 요청별 스타일 — ContextVar로 동시성 안전 (module global 잔존 금지).
_STYLE: ContextVar["ResolvedDeckStyle | None"] = ContextVar("pptx_style", default=None)
_SLIDE_META: ContextVar[tuple[int, int, str]] = ContextVar(
    "pptx_slide_meta", default=(1, 1, "cover"),
)


def _px_pt(px: float) -> float:
    return px * 0.75


def _hex_rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _colors():
    from pptx.dml.color import RGBColor

    style = _STYLE.get()
    if style is not None:
        t = style.template
        return {
            "accent": _hex_rgb(t.accent_color),
            "accent2": _hex_rgb(t.accent2),
            "blue": _hex_rgb(t.accent_color),
            "blue-light": _hex_rgb(t.accent2),
            "bg_dark": _hex_rgb(t.cover_bg or t.accent_color),
            "bg_card": _hex_rgb(t.surface_color),
            "orange": RGBColor(0xE8, 0x93, 0x24),
            "red": RGBColor(0xDC, 0x26, 0x26),
            "green": RGBColor(0x15, 0x80, 0x3D),
            "purple": RGBColor(0x7C, 0x3A, 0xED),
            "fg": _hex_rgb(t.body_color),
            "fg_muted": _hex_rgb(t.muted_color),
            "line": RGBColor(0xE5, 0xE7, 0xEB),
            "white": RGBColor(0xFF, 0xFF, 0xFF),
            "note_bg": RGBColor(0xFE, 0xF3, 0xC7),
            "note_fg": RGBColor(0x92, 0x40, 0x0E),
            "cover_box": _hex_rgb(t.surface_color),
            "cover_box_label": _hex_rgb(t.muted_color),
            "title": _hex_rgb(style.title_color),
            "page": _hex_rgb(t.page_number_color),
        }

    return {
        "accent": RGBColor(0x1A, 0x3A, 0x6B),
        "accent2": RGBColor(0x25, 0x63, 0xEB),
        "blue": RGBColor(0x1A, 0x3A, 0x6B),
        "blue-light": RGBColor(0x3B, 0x82, 0xF6),
        "bg_dark": RGBColor(0x1A, 0x3A, 0x6B),
        "bg_card": RGBColor(0xF8, 0xF9, 0xFB),
        "orange": RGBColor(0xE8, 0x93, 0x24),
        "red": RGBColor(0xDC, 0x26, 0x26),
        "green": RGBColor(0x15, 0x80, 0x3D),
        "purple": RGBColor(0x7C, 0x3A, 0xED),
        "fg": RGBColor(0x1A, 0x1A, 0x1A),
        "fg_muted": RGBColor(0x6B, 0x72, 0x80),
        "line": RGBColor(0xE5, 0xE7, 0xEB),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "note_bg": RGBColor(0xFE, 0xF3, 0xC7),
        "note_fg": RGBColor(0x92, 0x40, 0x0E),
        "cover_box": RGBColor(0x2B, 0x4A, 0x78),
        "cover_box_label": RGBColor(0xCF, 0xD6, 0xE4),
        "title": RGBColor(0xFF, 0xFF, 0xFF),
        "page": RGBColor(0x6B, 0x72, 0x80),
    }


def _color(name: str | None):
    c = _colors()
    return c.get(name or "blue", c["accent"])


def _new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_rect(slide, left, top, width, height, fill_color, *, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _add_textbox(
    slide, left, top, width, height, text, *,
    size=14, bold=False, color=None, align=None, anchor=None, wrap=True,
):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    colors = _colors()
    color = color or colors["fg"]
    align = align or PP_ALIGN.LEFT

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT
    run.font.color.rgb = color
    return box


def _add_multiline(
    slide, left, top, width, height, lines, *,
    size=12, bold=False, color=None, align=None, bullet=False, space_after=6,
):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    colors = _colors()
    color = color or colors["fg"]
    align = align or PP_ALIGN.LEFT

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = lines or []
    if not lines:
        lines = [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = f"•  {line}" if bullet else str(line)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = FONT
        run.font.color.rgb = color
    return box


def _add_header(slide, title, pageno, total_pages):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    colors = _colors()
    style = _STYLE.get()
    _add_rect(slide, 0, 0, SLIDE_W_IN, HEADER_H, colors["accent"])

    title_size = _px_pt(34)
    title_color = colors["white"]
    title_font = FONT
    if style is not None:
        title_size = style.title_size_pt
        title_color = colors["title"]
        title_font = style.title_font

    box = _add_textbox(
        slide, MARGIN, 0, SLIDE_W_IN - 2 * MARGIN - 1.6, HEADER_H, title,
        size=title_size, bold=True, color=title_color, anchor=MSO_ANCHOR.MIDDLE,
    )
    # _add_textbox uses FONT; override title font when styled
    if style is not None:
        for p in box.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = title_font
                run.font.size = Pt(title_size)

    # 스타일 모드에서는 하단 공통 페이지 번호를 쓰므로 헤더 번호는 생략.
    # 스타일 없음(8765)은 기존처럼 헤더에 번호 표시.
    if style is None:
        _add_textbox(
            slide, SLIDE_W_IN - MARGIN - 1.4, 0, 1.4, HEADER_H, f"{pageno}/{total_pages}",
            size=_px_pt(16), color=colors["white"], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )


def _add_page_number(slide, pageno: int, total_pages: int, pattern: str) -> None:
    """공통 페이지 번호. 표지 제외, pageno>=2, enabled일 때만."""
    from pptx.enum.text import PP_ALIGN

    style = _STYLE.get()
    if style is None or not style.page_number.enabled:
        return
    if pattern == "cover" or pageno < 2:
        return

    colors = _colors()
    pos = style.page_number.position
    text = f"{pageno}/{total_pages}"
    w, h = 1.6, 0.35
    top = FOOTER_TOP + 0.06
    if pos == "bottom-left":
        left, align = MARGIN, PP_ALIGN.LEFT
    elif pos == "bottom-center":
        left, align = (SLIDE_W_IN - w) / 2, PP_ALIGN.CENTER
    else:
        left, align = SLIDE_W_IN - MARGIN - w, PP_ALIGN.RIGHT
    _add_textbox(
        slide, left, top, w, h, text,
        size=_px_pt(14), color=colors["page"], align=align,
    )


def _add_subtitle(slide, subtitle):
    colors = _colors()
    _add_textbox(slide, MARGIN, CONTENT_TOP, CONTENT_W, 0.4, subtitle or "",
                 size=_px_pt(18), color=colors["fg_muted"])


def _add_footer(slide, company_name, deck_title, deck_date):
    from pptx.enum.text import PP_ALIGN

    colors = _colors()
    _add_rect(slide, MARGIN, FOOTER_TOP, CONTENT_W, 0.012, colors["line"])
    _add_textbox(
        slide, MARGIN, FOOTER_TOP + 0.06, CONTENT_W * 0.6, 0.35,
        f"{company_name} | {deck_title}", size=_px_pt(14), color=colors["fg_muted"],
    )
    _add_textbox(
        slide, MARGIN + CONTENT_W * 0.4, FOOTER_TOP + 0.06, CONTENT_W * 0.6, 0.35,
        f"Confidential | {deck_date}", size=_px_pt(14), color=colors["fg_muted"],
        align=PP_ALIGN.RIGHT,
    )


def _build_cover(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import PP_ALIGN

    colors = _colors()
    slide = _new_slide(prs)
    _add_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, colors["bg_dark"])

    left_w = SLIDE_W_IN * 0.55
    cy = SLIDE_H_IN / 2 - 1.7
    _add_textbox(slide, 1.05, cy, left_w, 0.55, d.get("company", ""),
                 size=_px_pt(36), bold=True, color=colors["orange"])
    style = _STYLE.get()
    cover_title_size = _px_pt(72)
    cover_title_color = colors["white"]
    cover_title_font = FONT
    if style is not None:
        cover_title_size = max(style.title_size_pt * 2, style.title_size_pt + 20)
        cover_title_color = colors["title"]
        cover_title_font = style.title_font
    title_box = _add_textbox(slide, 1.05, cy + 0.65, left_w, 2.0, d.get("title", ""),
                 size=cover_title_size, bold=True, color=cover_title_color)
    if style is not None:
        from pptx.util import Pt
        for p in title_box.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = cover_title_font
                run.font.size = Pt(cover_title_size)
    _add_rect(slide, 1.05, cy + 2.55, 1.25, 0.05, colors["orange"])
    _add_textbox(slide, 1.05, cy + 2.75, left_w, 0.6, d.get("subtitle", ""),
                 size=_px_pt(24), color=colors["white"])

    box_y = SLIDE_H_IN - 2.15
    box_w = (SLIDE_W_IN - 2 * 1.05 - 0.4) / 2
    for i, (label, value) in enumerate(
        [(_cover_labels[0], d.get("target", "")), (_cover_labels[1], d.get("date_version", ""))]
    ):
        bx = 1.05 + i * (box_w + 0.4)
        _add_rect(slide, bx, box_y, box_w, 1.0, colors["cover_box"])
        _add_textbox(slide, bx + 0.25, box_y + 0.12, box_w - 0.5, 0.3, label,
                     size=_px_pt(14), color=colors["cover_box_label"])
        _add_textbox(slide, bx + 0.25, box_y + 0.45, box_w - 0.5, 0.4, value,
                     size=_px_pt(20), bold=True, color=colors["white"])

    _add_textbox(
        slide, SLIDE_W_IN - 1.05 - 4.5, SLIDE_H_IN - 0.65, 4.5, 0.35,
        _cover_labels[2], size=_px_pt(14), color=colors["cover_box_label"],
        align=PP_ALIGN.RIGHT,
    )
    return slide


def _build_agenda(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    items = d.get("items", []) or []
    top = CONTENT_TOP + 0.55
    avail = CONTENT_BOTTOM - top
    row_h = min(0.85, avail / max(len(items), 1))
    y = top
    for it in items:
        _add_rect(slide, MARGIN, y + (row_h - 0.42) / 2, 0.9, 0.42, colors["accent"])
        _add_textbox(
            slide, MARGIN, y + (row_h - 0.42) / 2, 0.9, 0.42, str(it.get("ch", "")),
            size=_px_pt(16), bold=True, color=colors["white"], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_textbox(
            slide, MARGIN + 1.1, y, 8.5, row_h, it.get("title", ""),
            size=_px_pt(22), bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_textbox(
            slide, SLIDE_W_IN - MARGIN - 6.5, y, 6.5, row_h, it.get("summary", ""),
            size=_px_pt(18), color=colors["fg_muted"], align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_rect(slide, MARGIN, y + row_h - 0.012, CONTENT_W, 0.012, colors["line"])
        y += row_h

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _build_exec_summary(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import MSO_ANCHOR

    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    col_w = (CONTENT_W - 0.5) / 2
    top = CONTENT_TOP + 0.55
    cols = [
        (d.get("left_label", ""), d.get("left_items", []) or [], colors["red"]),
        (d.get("right_label", ""), d.get("right_items", []) or [], colors["green"]),
    ]
    for i, (label, items, hcolor) in enumerate(cols):
        x = MARGIN + i * (col_w + 0.5)
        _add_rect(slide, x, top, col_w, 0.55, hcolor)
        _add_textbox(slide, x + 0.2, top, col_w - 0.4, 0.55, label,
                     size=_px_pt(20), bold=True, color=colors["white"], anchor=MSO_ANCHOR.MIDDLE)
        body_top = top + 0.55
        body_h = CONTENT_BOTTOM - body_top
        _add_rect(slide, x, body_top, col_w, body_h, colors["bg_card"])
        iy = body_top + 0.2
        row_h = (body_h - 0.4) / max(len(items), 1)
        for it in items:
            _add_textbox(slide, x + 0.25, iy, col_w - 0.5, 0.35, it.get("title", ""),
                         size=_px_pt(20), bold=True)
            _add_textbox(slide, x + 0.25, iy + 0.38, col_w - 0.5, row_h - 0.4,
                         it.get("detail", ""), size=_px_pt(17), color=colors["fg_muted"])
            iy += row_h

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _build_metrics_row(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import PP_ALIGN

    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    metrics = d.get("metrics", []) or []
    n = max(len(metrics), 1)
    cols = min(n, 5)
    gap = 0.25
    w = (CONTENT_W - gap * (cols - 1)) / cols
    h = 2.2
    top = CONTENT_TOP + 0.6
    for i, m in enumerate(metrics):
        col, row = i % cols, i // cols
        x = MARGIN + col * (w + gap)
        y = top + row * (h + gap)
        _add_rect(slide, x, y, w, h, _color(m.get("color")))
        _add_textbox(slide, x, y + 0.4, w, 0.9, str(m.get("value", "")),
                     size=_px_pt(56), bold=True, color=colors["white"], align=PP_ALIGN.CENTER)
        _add_textbox(slide, x + 0.15, y + 1.35, w - 0.3, 0.7, m.get("label", ""),
                     size=_px_pt(18), color=colors["white"], align=PP_ALIGN.CENTER)

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _add_note(slide, text):
    colors = _colors()
    ny = CONTENT_BOTTOM - 0.5
    _add_rect(slide, MARGIN, ny, CONTENT_W, 0.5, colors["note_bg"])
    _add_textbox(slide, MARGIN + 0.25, ny, CONTENT_W - 0.5, 0.5, text,
                 size=_px_pt(18), color=colors["note_fg"])


def _build_compare_2col(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    note = d.get("footer_note")
    col_w = (CONTENT_W - 0.5) / 2
    top = CONTENT_TOP + 0.55
    body_top = top + 0.5
    body_bottom = CONTENT_BOTTOM - (0.65 if note else 0)
    body_h = body_bottom - body_top

    left_items = d.get("left_items", []) or []
    right_items = d.get("right_items", []) or []

    # left column: 3-col mini grid
    x = MARGIN
    _add_rect(slide, x, top, col_w, 0.5, colors["red"])
    _add_textbox(slide, x + 0.2, top, col_w - 0.4, 0.5, d.get("left_label", ""),
                 size=_px_pt(18), bold=True, color=colors["white"], anchor=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, x, body_top, col_w, body_h, colors["bg_card"])
    gcols = 3
    n = len(left_items)
    grows = max(math.ceil(n / gcols), 1)
    cw = (col_w - 0.4 - 0.15 * (gcols - 1)) / gcols
    ch = min(0.9, (body_h - 0.3 - 0.15 * (grows - 1)) / grows)
    for i, it in enumerate(left_items):
        gc, gr = i % gcols, i // gcols
        gx = x + 0.2 + gc * (cw + 0.15)
        gy = body_top + 0.15 + gr * (ch + 0.15)
        _add_rect(slide, gx, gy, cw, ch, colors["white"], line_color=colors["line"])
        _add_textbox(slide, gx + 0.05, gy, cw - 0.1, ch, str(it),
                     size=_px_pt(18), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # right column: highlighted first item + list
    x2 = MARGIN + col_w + 0.5
    _add_rect(slide, x2, top, col_w, 0.5, colors["green"])
    _add_textbox(slide, x2 + 0.2, top, col_w - 0.4, 0.5, d.get("right_label", ""),
                 size=_px_pt(18), bold=True, color=colors["white"], anchor=MSO_ANCHOR.MIDDLE)
    _add_rect(slide, x2, body_top, col_w, body_h, colors["bg_card"])
    ry = body_top + 0.15
    if right_items:
        _add_rect(slide, x2 + 0.2, ry, col_w - 0.4, 0.8, colors["accent"])
        _add_textbox(slide, x2 + 0.3, ry, col_w - 0.6, 0.8, str(right_items[0]),
                     size=_px_pt(22), bold=True, color=colors["white"], align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        ry += 0.95
        for it in right_items[1:]:
            _add_rect(slide, x2 + 0.2, ry, col_w - 0.4, 0.55, colors["white"], line_color=colors["line"])
            _add_textbox(slide, x2 + 0.3, ry, col_w - 0.6, 0.55, str(it),
                         size=_px_pt(18), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            ry += 0.65

    if note:
        _add_note(slide, note)

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _build_card_grid_4(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    from src.engine.output.deck.card_content_guard import coerce_intro, coerce_outro

    y = CONTENT_TOP + 0.5
    intro = coerce_intro(d.get("intro"))
    if intro:
        lines = intro.get("lines", []) or []
        label = intro.get("label", "") or ""
        # label만 있으면 1줄, lines만 있으면 lines 높이 (빈 textbox 방지)
        h = 0.28 + (0.28 if label else 0) + 0.28 * max(len(lines), 0)
        h = max(h, 0.45)
        _add_rect(slide, MARGIN, y, CONTENT_W, h, colors["bg_card"])
        ty = y + 0.08
        if label:
            _add_textbox(slide, MARGIN + 0.2, ty, CONTENT_W - 0.4, 0.28, label,
                         size=_px_pt(20), bold=True, color=colors["accent"])
            ty += 0.3
        if lines:
            _add_multiline(slide, MARGIN + 0.2, ty, CONTENT_W - 0.4, h - (ty - y),
                           lines, size=_px_pt(15))
        y += h + 0.25

    section_title = (d.get("section_title") or "").strip()
    if section_title:
        _add_textbox(slide, MARGIN, y, CONTENT_W, 0.35, section_title,
                     size=_px_pt(20), bold=True, color=colors["accent"])
        y += 0.45

    cards = d.get("cards", []) or []
    outro = coerce_outro(d.get("outro"))
    if outro:
        bullets = outro.get("bullets", []) or []
        olabel = outro.get("label", "") or ""
        outro_h = 0.28 + (0.28 if olabel else 0) + 0.28 * max(len(bullets), 0)
        outro_h = max(outro_h, 0.45)
    else:
        outro_h = 0
    grid_bottom = CONTENT_BOTTOM - (outro_h + 0.25 if outro else 0)
    grid_h = grid_bottom - y
    cols = min(len(cards), 4) or 1
    gap = 0.2
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    for i, c in enumerate(cards):
        x = MARGIN + i * (cw + gap)
        _add_rect(slide, x, y, cw, 0.65, _color(c.get("color")))
        _add_textbox(slide, x + 0.1, y, cw - 0.2, 0.65, c.get("title", ""),
                     size=_px_pt(18), bold=True, color=colors["white"], align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(slide, x, y + 0.65, cw, grid_h - 0.65, colors["white"], line_color=colors["line"])
        _add_textbox(slide, x + 0.12, y + 0.78, cw - 0.24, grid_h - 0.9, c.get("subtitle", ""),
                     size=_px_pt(15), color=colors["fg_muted"])

    if outro:
        bullets = outro.get("bullets", []) or []
        olabel = outro.get("label", "") or ""
        oy = CONTENT_BOTTOM - outro_h
        _add_rect(slide, MARGIN, oy, CONTENT_W, outro_h, colors["bg_card"])
        ty = oy + 0.08
        if olabel:
            _add_textbox(slide, MARGIN + 0.2, ty, CONTENT_W - 0.4, 0.28, olabel,
                         size=_px_pt(20), bold=True, color=colors["accent"])
            ty += 0.3
        if bullets:
            _add_multiline(slide, MARGIN + 0.2, ty, CONTENT_W - 0.4, outro_h - (ty - oy),
                           bullets, size=_px_pt(15), bullet=True)

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _build_phase_roadmap(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    phases = d.get("phases", []) or []
    note = d.get("footer_note")
    top = CONTENT_TOP + 0.55
    bottom = CONTENT_BOTTOM - (0.65 if note else 0)
    n = max(len(phases), 1)
    gap = 0.2
    w = (CONTENT_W - gap * (n - 1)) / n
    for i, ph in enumerate(phases):
        x = MARGIN + i * (w + gap)
        _add_rect(slide, x, top, w, 0.55, _color(ph.get("color")))
        _add_textbox(slide, x + 0.08, top, w - 0.16, 0.55, ph.get("label", ""),
                     size=_px_pt(16), bold=True, color=colors["white"], align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        body_top = top + 0.55
        _add_rect(slide, x, body_top, w, bottom - body_top, colors["bg_card"])
        ty = body_top + 0.15
        _add_textbox(slide, x + 0.15, ty, w - 0.3, 0.4, ph.get("title", ""),
                     size=_px_pt(16), bold=True)
        ty += 0.4
        _add_textbox(slide, x + 0.15, ty, w - 0.3, 0.3, ph.get("period", ""),
                     size=_px_pt(13), color=colors["fg_muted"])
        ty += 0.35
        _add_multiline(slide, x + 0.15, ty, w - 0.3, bottom - ty - 0.1, ph.get("tasks", []),
                       size=_px_pt(14))

    if note:
        _add_note(slide, note)

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _build_dimension_5(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    dims = d.get("dimensions", []) or []
    note = d.get("footer_note")
    top = CONTENT_TOP + 0.55
    bottom = CONTENT_BOTTOM - (0.65 if note else 0)
    n = max(len(dims), 1)
    gap = 0.18
    w = (CONTENT_W - gap * (n - 1)) / n
    for i, dim in enumerate(dims):
        x = MARGIN + i * (w + gap)
        c = _color(dim.get("color"))
        _add_rect(slide, x, top, w, 0.55, c)
        _add_textbox(slide, x + 0.08, top, w - 0.16, 0.55, dim.get("label", ""),
                     size=_px_pt(15), bold=True, color=colors["white"], align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
        body_top = top + 0.55
        _add_rect(slide, x, body_top, w, bottom - body_top, colors["bg_card"])
        _add_multiline(slide, x + 0.15, body_top + 0.15, w - 0.3, bottom - body_top - 0.3,
                       dim.get("bullets", []), size=_px_pt(14), bullet=True)

    if note:
        _add_note(slide, note)

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _build_table(prs, d, deck, pageno, total_pages):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    columns = d.get("columns") or []
    rows = d.get("rows") or []
    ncols = max(len(columns), 1)
    nrows = len(rows)
    top = CONTENT_TOP + 0.55
    bottom = CONTENT_BOTTOM - (0.55 if d.get("footnote") else 0)
    avail_h = bottom - top
    row_h = min(0.55, avail_h / max(nrows + 1, 1))
    col_w = CONTENT_W / ncols
    hdr_h = min(0.45, row_h)

    for ci, col in enumerate(columns):
        x = MARGIN + ci * col_w
        _add_rect(slide, x, top, col_w - 0.05, hdr_h, colors["accent"])
        _add_textbox(
            slide, x + 0.08, top, col_w - 0.16, hdr_h,
            str(col.get("label", "")),
            size=_px_pt(14), bold=True, color=colors["white"],
            anchor=MSO_ANCHOR.MIDDLE,
        )

    y = top + hdr_h
    for ri, row in enumerate(rows):
        for ci, col in enumerate(columns):
            x = MARGIN + ci * col_w
            key = col.get("key", "")
            val = str(row.get(key, "")) if isinstance(row, dict) else ""
            bg = colors["bg_card"] if ri % 2 else colors["white"]
            _add_rect(slide, x, y, col_w - 0.05, row_h, bg)
            _add_textbox(
                slide, x + 0.08, y + 0.04, col_w - 0.16, row_h - 0.08, val,
                size=_px_pt(13), anchor=MSO_ANCHOR.TOP,
            )
        y += row_h

    if d.get("footnote"):
        _add_note(slide, str(d.get("footnote")))

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _add_key_message(slide, text, top):
    """서술형/요약형 공통 — 핵심 메시지 콜아웃 박스. 반환: 다음 y."""
    colors = _colors()
    if not text:
        return top
    box_h = 0.9
    _add_rect(slide, MARGIN, top, 0.09, box_h, colors["accent"])
    _add_rect(slide, MARGIN + 0.09, top, CONTENT_W - 0.09, box_h, colors["bg_card"])
    from pptx.enum.text import MSO_ANCHOR
    _add_textbox(
        slide, MARGIN + 0.42, top, CONTENT_W - 0.7, box_h, str(text),
        size=_px_pt(24), bold=True, anchor=MSO_ANCHOR.MIDDLE,
    )
    return top + box_h + 0.4


def _build_narrative(prs, d, deck, pageno, total_pages):
    """서술형(narrative) — 핵심 메시지 + 문장형 단락."""
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    top = CONTENT_TOP + 0.55
    top = _add_key_message(slide, d.get("key_message"), top)

    paragraphs = [str(p) for p in (d.get("paragraphs") or []) if str(p).strip()]
    _add_multiline(
        slide, MARGIN, top, CONTENT_W, CONTENT_BOTTOM - top, paragraphs,
        size=_px_pt(23), space_after=16,
    )

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


def _build_summary(prs, d, deck, pageno, total_pages):
    """요약형(summary) — 핵심 메시지 + 짧은 불릿."""
    colors = _colors()
    slide = _new_slide(prs)
    _add_header(slide, d.get("title", deck.title), pageno, total_pages)
    _add_subtitle(slide, d.get("subtitle", ""))

    top = CONTENT_TOP + 0.55
    top = _add_key_message(slide, d.get("key_message"), top)

    points = [str(p) for p in (d.get("points") or []) if str(p).strip()]
    _add_multiline(
        slide, MARGIN, top, CONTENT_W, CONTENT_BOTTOM - top, points,
        size=_px_pt(24), bullet=True, space_after=14, color=colors["fg"],
    )

    _add_footer(slide, deck.company_name, deck.title, deck.date)
    return slide


_PATTERN_BUILDERS = {
    "cover": _build_cover,
    "agenda": _build_agenda,
    "exec-summary": _build_exec_summary,
    "metrics-row": _build_metrics_row,
    "compare-2col": _build_compare_2col,
    "card-grid-4": _build_card_grid_4,
    "phase-roadmap": _build_phase_roadmap,
    "dimension-5": _build_dimension_5,
    "table": _build_table,
    "narrative": _build_narrative,
    "summary": _build_summary,
}


def render_deck_to_pptx(
    deck: Deck, out_path: Path | None = None,
    *, on_progress: "ProgressFn | None" = None,
    template_id: str | None = None,
    master_style: dict | None = None,
    page_number: dict | None = None,
    style: "ResolvedDeckStyle | None" = None,
) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    if out_path is None:
        out_path = Path(tempfile.mkdtemp(prefix="iris_pptx_")) / "deck.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if style is None and (template_id or master_style or page_number):
        from src.engine.output.deck.theme import resolve_deck_style
        style = resolve_deck_style(
            template_id or "clean-light",
            master_style=master_style,
            page_number=page_number,
        )

    # 슬라이드 전체가 같은 폰트/커버 라벨 언어를 쓰도록 첫 판단을 한 번만
    # 하고 모듈 전역에 담아둠 — 콘텐츠 언어(소스 문서 언어)를 따라감.
    global FONT, _cover_labels
    lang = detect_deck_lang(deck)
    FONT = _FONT_BY_LANG[lang]
    _cover_labels = _COVER_LABELS[lang]

    token = _STYLE.set(style)
    try:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)

        total = len(deck.slides)
        for i, slide in enumerate(deck.slides, 1):
            builder = _PATTERN_BUILDERS.get(slide.pattern)
            if builder is None:
                raise ValueError(f"미지 패턴: {slide.pattern}")
            built = builder(prs, slide.data, deck, i, total)
            _add_page_number(built, i, total, slide.pattern)
            if on_progress is not None:
                label = str(slide.data.get("title", slide.data.get("company", slide.pattern)))[:40]
                on_progress(i, total, label)

        prs.save(str(out_path))
        return out_path
    finally:
        _STYLE.reset(token)


__all__ = ["render_deck_to_pptx"]
